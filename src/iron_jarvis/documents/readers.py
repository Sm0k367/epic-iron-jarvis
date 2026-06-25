"""Document text extraction (readers).

``extract_text(path)`` reads a real-world file and returns its text content,
dispatching on the lowercased filename suffix:

* ``.pdf``  -> pypdf (join every page's ``extract_text()``).
* ``.docx`` -> python-docx (paragraph text + table cell text).
* ``.xlsx`` -> openpyxl (``read_only``/``data_only``; "## <sheet>" + TSV rows).
* ``.pptx`` -> python-pptx (per slide, join each shape's text).
* ``.csv``  -> stdlib csv (tab-joined rows).
* ``.txt/.md/.json/.py/.js/.ts/.html/.yaml/.yml/.log`` and any unknown-but-text
  file -> decoded as UTF-8 (``errors="replace"``).
* ``.png/.jpg/.jpeg/.gif/.bmp/.webp`` -> Pillow, returning a concise note such as
  ``"[image PNG 800x600, mode RGB]"`` (NO OCR).

A clear :class:`ValueError` is raised for a genuinely unsupported / binary type.
"""

from __future__ import annotations

import csv
from pathlib import Path

#: Suffixes read verbatim as UTF-8 text.
_TEXT_SUFFIXES: frozenset[str] = frozenset(
    {
        ".txt",
        ".md",
        ".json",
        ".py",
        ".js",
        ".ts",
        ".html",
        ".htm",
        ".yaml",
        ".yml",
        ".log",
    }
)

#: Raster image suffixes -> described, never OCR'd.
_IMAGE_SUFFIXES: frozenset[str] = frozenset(
    {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}
)

#: Structured document suffixes handled by dedicated parsers.
_DOC_SUFFIXES: frozenset[str] = frozenset({".pdf", ".docx", ".xlsx", ".pptx", ".csv"})

#: Every suffix ``extract_text`` knows how to read (unknown text files also work,
#: but are not advertised here).
SUPPORTED_READ: set[str] = set(_DOC_SUFFIXES | _TEXT_SUFFIXES | _IMAGE_SUFFIXES)


def extract_text(path: str | Path) -> str:
    """Return the text content of ``path``, dispatched by file suffix.

    Raises :class:`ValueError` for a truly unsupported / binary file type and
    :class:`FileNotFoundError` if the path does not exist.
    """
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"no such file: {p}")
    if p.is_dir():
        raise ValueError(f"path is a directory, not a document: {p}")

    suffix = p.suffix.lower()
    if suffix == ".pdf":
        return _read_pdf(p)
    if suffix == ".docx":
        return _read_docx(p)
    if suffix == ".xlsx":
        return _read_xlsx(p)
    if suffix == ".pptx":
        return _read_pptx(p)
    if suffix == ".csv":
        return _read_csv(p)
    if suffix in _IMAGE_SUFFIXES:
        return _describe_image(p)
    if suffix in _TEXT_SUFFIXES:
        return p.read_text(encoding="utf-8", errors="replace")

    # Unknown suffix: sniff for binary content; decode if it looks like text.
    data = p.read_bytes()
    if b"\x00" in data:
        raise ValueError(
            f"unsupported binary file type: {suffix or p.name!r}"
        )
    return data.decode("utf-8", errors="replace")


# --- structured-format readers ------------------------------------------------


def _read_pdf(p: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(p))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _read_docx(p: Path) -> str:
    import docx

    doc = docx.Document(str(p))
    parts: list[str] = [para.text for para in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _read_xlsx(p: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(filename=str(p), read_only=True, data_only=True)
    try:
        parts: list[str] = []
        for ws in wb.worksheets:
            parts.append(f"## {ws.title}")
            for row in ws.iter_rows(values_only=True):
                parts.append("\t".join("" if c is None else str(c) for c in row))
        return "\n".join(parts)
    finally:
        wb.close()


def _read_pptx(p: Path) -> str:
    import pptx

    prs = pptx.Presentation(str(p))
    parts: list[str] = []
    for slide in prs.slides:
        texts = [
            shape.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text
        ]
        parts.append("\n".join(texts))
    return "\n".join(parts)


def _read_csv(p: Path) -> str:
    rows: list[str] = []
    with open(p, newline="", encoding="utf-8", errors="replace") as f:
        for row in csv.reader(f):
            rows.append("\t".join(row))
    return "\n".join(rows)


def _describe_image(p: Path) -> str:
    from PIL import Image

    with Image.open(p) as img:
        fmt = img.format or (p.suffix.lstrip(".").upper() or "IMAGE")
        width, height = img.size
        mode = img.mode
    return f"[image {fmt} {width}x{height}, mode {mode}]"
